from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.graphics.shapes import Drawing, Rect, String
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw, ImageFont

# 1. Create simple temporary logos for Style 1 and Style 2
def create_logo(text, filename, bg_color, text_color):
    img = Image.new("RGB", (600, 200), bg_color)
    draw = ImageDraw.Draw(img)
    font = ImageFont.load_default()
    w, h = draw.textsize(text, font=font)
    draw.text(((600 - w) / 2, (200 - h) / 2), text, fill=text_color, font=font)
    img.save(f"/mnt/data/{filename}")

create_logo("GHANI TECHNO TRADING COMPANY (Style 1)", "logo_style1.png", (10, 40, 120), (255, 255, 255))
create_logo("GHANI TECHNO TRADING COMPANY (Style 2)", "logo_style2.png", (0, 0, 0), (255, 215, 0))

# 2. PDF Company Profile
pdf_path = "/mnt/data/company_profile.pdf"
styles = getSampleStyleSheet()
story = []

story.append(Paragraph("Ghani Techno Trading Company<br/><br/>Company Profile", styles['Title']))
story.append(Spacer(1, 12))
story.append(Paragraph("We provide engineering, industrial, medical, construction, and general trading solutions with a commitment to reliability, professionalism, and long-term partnerships.", styles['BodyText']))
story.append(Spacer(1, 12))

doc = SimpleDocTemplate(pdf_path, pagesize=A4)
doc.build(story)

# 3. PowerPoint Presentation
ppt_path = "/mnt/data/company_intro.pptx"
prs = Presentation()
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Ghani Techno Trading Company"
subtitle.text = "Engineering • Industrial • Medical • Construction • General Trading"

prs.save(ppt_path)

# 4. Organizational Chart PNG
img = Image.new("RGB", (900, 600), "white")
draw = ImageDraw.Draw(img)

# Basic org chart layout
draw.rectangle((300, 50, 600, 120), outline="black")
draw.text((330, 80), "CEO", fill="black")

# Divisions
positions = ["Hardware & Firmware", "Software Development", "Mobile Apps",
             "Cybersecurity", "Marketing & Creative", "R&D & Innovation",
             "Data & Analytics", "HR & Admin"]

y = 180
for pos in positions:
    draw.rectangle((100, y, 800, y+50), outline="black")
    draw.text((120, y+15), pos, fill="black")
    y += 60

orgchart_path = "/mnt/data/org_chart.png"
img.save(orgchart_path)

# 5. Canva Brochure (PDF)
brochure_path = "/mnt/data/canva_brochure.pdf"
bro_doc = SimpleDocTemplate(brochure_path, pagesize=A4)
bro_story = []

bro_story.append(Paragraph("Ghani Techno Trading Company<br/><br/>Brochure", styles['Title']))
bro_story.append(Spacer(1, 12))
bro_story.append(Paragraph("A multi-domain trading and engineering company dealing in electrical, mechanical, hydraulic, optronics, medical equipment, construction materials, and general order supplies.", styles['BodyText']))

bro_doc.build(bro_story)

pdf_path, ppt_path, brochure_path, orgchart_path, "/mnt/data/logo_style1.png", "/mnt/data/logo_style2.png"
